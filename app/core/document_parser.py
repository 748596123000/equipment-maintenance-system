"""
多格式文档解析模块

支持解析以下文件格式：
- PDF (.pdf) — 使用 PyMuPDF + pdfplumber
- Word (.docx) — 使用 python-docx
- Excel (.xlsx) — 使用 openpyxl
- PowerPoint (.pptx) — 使用 python-pptx
- 文本 (.txt, .md, .csv, .json, .xml, .log) — 直接读取
- 图片 (.jpg, .jpeg, .png, .bmp, .gif, .tiff, .webp) — 调用多模态模型生成描述

所有解析器返回统一格式，与 PDFParser.parse_pdf() 的输出兼容。
"""

import base64
import json
import logging
import os
from typing import Dict, List, Optional

from app.config import settings

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "xlsx", "pptx",
    "txt", "md", "csv", "json", "xml", "log",
    "jpg", "jpeg", "png", "bmp", "gif", "tiff", "webp",
}

MIME_TYPE_MAP = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "txt": "text/plain",
    "md": "text/markdown",
    "csv": "text/csv",
    "json": "application/json",
    "xml": "application/xml",
    "log": "text/plain",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "png": "image/png",
    "bmp": "image/bmp",
    "gif": "image/gif",
    "tiff": "image/tiff",
    "webp": "image/webp",
}

MAGIC_BYTES_MAP = {
    "pdf": b"%PDF",
    "docx": b"PK\x03\x04",
    "xlsx": b"PK\x03\x04",
    "pptx": b"PK\x03\x04",
    "png": b"\x89PNG",
    "jpg": b"\xff\xd8\xff",
    "jpeg": b"\xff\xd8\xff",
    "gif": b"GIF8",
    "bmp": b"BM",
    "webp": b"RIFF",
}


def get_file_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def validate_file_type(filename: str, content: bytes) -> Optional[str]:
    ext = get_file_extension(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        return None

    if ext in MAGIC_BYTES_MAP:
        expected = MAGIC_BYTES_MAP[ext]
        if not content.startswith(expected):
            if ext in ("docx", "xlsx", "pptx"):
                if not content.startswith(b"PK"):
                    return None
            else:
                return None

    return ext


def parse_document(file_path: str, filename: Optional[str] = None, max_pages: Optional[int] = None) -> dict:
    if filename is None:
        filename = os.path.basename(file_path)

    ext = get_file_extension(filename)

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    parsers = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "xlsx": _parse_xlsx,
        "pptx": _parse_pptx,
    }

    text_extensions = {"txt", "md", "csv", "json", "xml", "log"}
    image_extensions = {"jpg", "jpeg", "png", "bmp", "gif", "tiff", "webp"}

    if ext in parsers:
        return parsers[ext](file_path, filename, max_pages=max_pages)
    elif ext in text_extensions:
        return _parse_text(file_path, filename, ext)
    elif ext in image_extensions:
        return _parse_image(file_path, filename, ext)
    else:
        raise ValueError(f"不支持的文件格式: .{ext}")


def _parse_pdf(file_path: str, filename: str, max_pages: Optional[int] = None) -> dict:
    from app.core.pdf_parser import PDFParser
    parser = PDFParser(file_path, max_pages=max_pages)
    return parser.parse_pdf()


def _parse_docx(file_path: str, filename: str, max_pages: Optional[int] = None) -> dict:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("需要安装 python-docx 库: pip install python-docx")

    doc = Document(file_path)
    paragraphs = []
    para_index = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""
        if "Heading 1" in style_name or "标题 1" in style_name:
            level = 1
        elif "Heading 2" in style_name or "标题 2" in style_name:
            level = 2
        elif "Heading 3" in style_name or "标题 3" in style_name:
            level = 3
        else:
            level = 0

        paragraphs.append({
            "page": 1,
            "title": text if level > 0 else "",
            "content": text,
            "level": level,
        })
        para_index += 1

    table_index = 0
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_rows.append(row_data)

        if table_rows:
            paragraphs.append({
                "page": 1,
                "title": f"[表格 {table_index + 1}]",
                "content": _format_table_text(table_rows),
                "level": 0,
            })
            table_index += 1

    return _build_result(filename, paragraphs, doc_type="docx")


def _parse_xlsx(file_path: str, filename: str, max_pages: Optional[int] = None) -> dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("需要安装 openpyxl 库: pip install openpyxl")

    wb = load_workbook(file_path, read_only=True, data_only=True)
    paragraphs = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]

        paragraphs.append({
            "page": sheet_idx + 1,
            "title": f"工作表: {sheet_name}",
            "content": f"工作表: {sheet_name}",
            "level": 1,
        })

        table_rows = []
        for row in ws.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in row_data):
                table_rows.append(row_data)

        if table_rows:
            paragraphs.append({
                "page": sheet_idx + 1,
                "title": f"[数据表 {sheet_idx + 1}]",
                "content": _format_table_text(table_rows),
                "level": 0,
            })

    wb.close()
    return _build_result(filename, paragraphs, doc_type="xlsx", page_count=len(wb.sheetnames))


def _parse_pptx(file_path: str, filename: str, max_pages: Optional[int] = None) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("需要安装 python-pptx 库: pip install python-pptx")

    prs = Presentation(file_path)
    paragraphs = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        paragraphs.append({
            "page": slide_num,
            "title": f"幻灯片 {slide_num}",
            "content": f"幻灯片 {slide_num}",
            "level": 1,
        })

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs.append({
                            "page": slide_num,
                            "title": "",
                            "content": text,
                            "level": 0,
                        })

            if shape.has_table:
                table_rows = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_rows.append(row_data)
                if table_rows:
                    paragraphs.append({
                        "page": slide_num,
                        "title": f"[表格]",
                        "content": _format_table_text(table_rows),
                        "level": 0,
                    })

    return _build_result(filename, paragraphs, doc_type="pptx", page_count=len(prs.slides))


def _parse_text(file_path: str, filename: str, ext: str) -> dict:
    encodings = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]
    content = None

    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read()
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    if content is None:
        raise RuntimeError(f"无法解码文本文件: {filename}")

    paragraphs = []
    if ext == "json":
        paragraphs = _parse_json_content(content, filename)
    elif ext == "csv":
        paragraphs = _parse_csv_content(content, filename)
    elif ext == "xml":
        paragraphs.append({
            "page": 1,
            "title": "",
            "content": content,
            "level": 0,
        })
    else:
        lines = content.split("\n")
        current_parts = []
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if current_parts:
                    paragraphs.append({
                        "page": 1,
                        "title": "",
                        "content": "\n".join(current_parts),
                        "level": 0,
                    })
                    current_parts = []
            else:
                current_parts.append(stripped)

        if current_parts:
            paragraphs.append({
                "page": 1,
                "title": "",
                "content": "\n".join(current_parts),
                "level": 0,
            })

    return _build_result(filename, paragraphs, doc_type=ext)


def _parse_json_content(content: str, filename: str) -> List[Dict]:
    paragraphs = []
    try:
        data = json.loads(content)
        formatted = json.dumps(data, ensure_ascii=False, indent=2)
        paragraphs.append({
            "page": 1,
            "title": "",
            "content": formatted,
            "level": 0,
        })
    except json.JSONDecodeError:
        paragraphs.append({
            "page": 1,
            "title": "",
            "content": content,
            "level": 0,
        })
    return paragraphs


def _parse_csv_content(content: str, filename: str) -> List[Dict]:
    import csv
    import io

    paragraphs = []
    reader = csv.reader(io.StringIO(content))
    table_rows = []
    for row in reader:
        cleaned = [cell.strip() for cell in row]
        if any(cell for cell in cleaned):
            table_rows.append(cleaned)

    if table_rows:
        paragraphs.append({
            "page": 1,
            "title": "[CSV数据]",
            "content": _format_table_text(table_rows),
            "level": 0,
        })
    return paragraphs


def _parse_image(file_path: str, filename: str, ext: str) -> dict:
    with open(file_path, "rb") as f:
        image_bytes = f.read()

    paragraphs = []

    ocr_text = _ocr_image(image_bytes)
    if ocr_text:
        paragraphs.append({
            "page": 1,
            "title": "[OCR识别文字]",
            "content": f"[OCR识别文字] {ocr_text}",
            "level": 0,
        })

    description = _describe_image_file(image_bytes, ext)
    if description:
        paragraphs.append({
            "page": 1,
            "title": "[图片描述]",
            "content": f"[图片描述] {description}",
            "level": 0,
        })

    if not paragraphs:
        paragraphs.append({
            "page": 1,
            "title": "",
            "content": f"[图片文件: {filename}]",
            "level": 0,
        })

    return _build_result(filename, paragraphs, doc_type=ext, page_count=1)


def _ocr_image(image_bytes: bytes) -> Optional[str]:
    if settings.OCR_BACKEND == "none":
        return None

    try:
        from app.services.ocr_service import get_ocr_service
        ocr = get_ocr_service()
        result = ocr.ocr_image(image_bytes)
        if result and result.text and result.confidence > 0.3:
            return result.text
    except Exception as e:
        logger.warning(f"OCR 处理失败: {e}")

    return None


def _describe_image_file(image_bytes: bytes, ext: str) -> Optional[str]:
    try:
        from app.services.vision_service import get_vision_service
        vision = get_vision_service()
        return vision.describe_image(image_bytes, ext)
    except Exception as e:
        logger.warning(f"视觉模型调用失败: {e}")

    if not settings.dashscope_api_key:
        logger.warning("未配置DASHSCOPE_API_KEY，跳过图片描述生成")
        return None

    try:
        from openai import OpenAI

        client = OpenAI(
            api_key=settings.dashscope_api_key,
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1"
        )

        mime_map = {
            "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "png": "image/png", "bmp": "image/bmp",
            "gif": "image/gif", "tiff": "image/tiff",
            "webp": "image/webp",
        }
        mime_type = mime_map.get(ext, "image/png")
        base64_str = base64.b64encode(image_bytes).decode("utf-8")

        response = client.chat.completions.create(
            model="qwen-vl-max",
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime_type};base64,{base64_str}"}
                    },
                    {
                        "type": "text",
                        "text": "请详细描述这张设备检修相关图片的内容，包括设备名称、部件、操作步骤等"
                    }
                ]
            }],
            max_tokens=300
        )

        return response.choices[0].message.content

    except Exception as e:
        logger.warning(f"调用通义千问多模态模型失败: {e}")
        return None


def _format_table_text(rows: List[List[str]]) -> str:
    if not rows:
        return ""

    max_cols = max(len(row) for row in rows)
    normalized = []
    for row in rows:
        padded = row + [""] * (max_cols - len(row))
        normalized.append(padded)

    col_widths = []
    for col_idx in range(max_cols):
        width = max(len(row[col_idx]) for row in normalized)
        col_widths.append(min(width, 50))

    lines = []
    for row in normalized:
        parts = []
        for col_idx, cell in enumerate(row):
            parts.append(cell.ljust(col_widths[col_idx]))
        lines.append(" | ".join(parts))

    return "\n".join(lines)


def _build_result(
    filename: str,
    paragraphs: List[Dict],
    doc_type: str = "unknown",
    page_count: int = 1,
) -> dict:
    return {
        "filename": filename,
        "metadata": {
            "doc_type": doc_type,
            "total_pages": page_count,
        },
        "paragraphs": paragraphs,
        "tables": [],
        "images": [],
        "total_pages": page_count,
        "paragraph_count": len(paragraphs),
        "table_count": 0,
        "image_count": 0,
    }
